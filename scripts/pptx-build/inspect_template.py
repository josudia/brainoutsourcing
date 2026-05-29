#!/usr/bin/env python3
"""Inspect a PPTX template: layouts, placeholders, existing slides, slide size."""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def main():
    if len(sys.argv) < 2:
        sys.exit("Usage: inspect_template.py <template.pptx>")
    path = Path(sys.argv[1])
    if not path.exists():
        sys.exit(f"Not found: {path}")

    prs = Presentation(str(path))

    print("=== SLIDE LAYOUTS ===")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"{i}: {layout.name}")
        for ph in layout.placeholders:
            try:
                pf = ph.placeholder_format
                print(f"    idx={pf.idx} type={pf.type} name={ph.name}")
            except Exception as e:
                print(f"    (placeholder read error: {e})")

    print()
    print("=== EXISTING SLIDES ===")
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}: layout={slide.slide_layout.name}")

    print()
    print("=== SLIDE SIZE ===")
    print(
        f"width={prs.slide_width} ({Emu(prs.slide_width).inches:.2f}in) "
        f"height={prs.slide_height} ({Emu(prs.slide_height).inches:.2f}in)"
    )


if __name__ == "__main__":
    main()
