#!/usr/bin/env python3
"""Inspect a built PPTX: per-slide title, body text, notes, layout, placeholders."""
import sys
from pathlib import Path

from pptx import Presentation


def main():
    if len(sys.argv) < 2:
        sys.exit("Usage: inspect_output.py <pptx>")
    path = Path(sys.argv[1])
    if not path.exists():
        sys.exit(f"Not found: {path}")

    prs = Presentation(str(path))
    print(f"=== {path.name} ===")
    print(f"Slides: {len(prs.slides)}")
    print()

    for i, slide in enumerate(prs.slides, 1):
        print(f"--- Slide {i} (layout='{slide.slide_layout.name}') ---")
        for ph in slide.placeholders:
            try:
                pf = ph.placeholder_format
                tag = f"ph idx={pf.idx} type={pf.type} name={ph.name}"
            except Exception:
                tag = f"ph (n/a) name={ph.name}"
            if ph.has_text_frame:
                txt = ph.text_frame.text
                txt_preview = (txt[:200] + "...") if len(txt) > 200 else txt
                txt_preview = txt_preview.replace("\n", "\\n")
                print(f"  {tag}: '{txt_preview}'")
            else:
                print(f"  {tag}: (no text frame)")
        # Non-placeholder shapes (e.g. textboxes inserted as fallback)
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            if shape.has_text_frame:
                txt = shape.text_frame.text
                txt_preview = (txt[:200] + "...") if len(txt) > 200 else txt
                txt_preview = txt_preview.replace("\n", "\\n")
                print(f"  shape '{shape.name}': '{txt_preview}'")
            else:
                print(f"  shape '{shape.name}': (no text)")
        # Notes
        if slide.has_notes_slide:
            notes_txt = slide.notes_slide.notes_text_frame.text
            notes_preview = (notes_txt[:300] + "...") if len(notes_txt) > 300 else notes_txt
            notes_preview = notes_preview.replace("\n", "\\n")
            print(f"  NOTES: '{notes_preview}'")
        print()


if __name__ == "__main__":
    main()
