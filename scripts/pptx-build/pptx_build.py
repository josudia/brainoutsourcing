#!/usr/bin/env python3
"""
pptx_build.py — generic PowerPoint builder.

Reads a Markdown slide file (see README for the convention) and builds a new
PPTX from a template PPTX (master/layouts/fonts/colors are preserved). Content
is written into the first matching placeholders. Speaker notes are optional via
`:::notes:::` blocks.

Usage:
  ./venv/bin/python pptx_build.py <markdown> --template <template.pptx> --output <new.pptx>

Markdown convention:
  ## Slide 1: Slide title
  Optional subtitle or first paragraph.

  - Bullet 1
  - Bullet 2
    - Sub-bullet
  - Bullet 3

  :::notes:::
  Speaker notes for this slide. Multiple lines allowed.
  :::endnotes:::

  ## Slide 2 (layout=2): Another slide with an explicit layout index
  ...

Layout index: 0 = Title, 1 = Title + Content (default), 2 = Section Header,
              5 = Title Only, 6 = Blank. Depends on the template master.

Image slots: kept in the template as placeholders and marked with
             "[IMAGE: <description from the Markdown :::image::: block>]".
"""

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


SLIDE_HEADING_RE = re.compile(
    r"^##\s+Slide\s+\d+\s*(?:\(layout=(\d+)\))?\s*:?\s*(.+?)\s*$"
)
NOTES_START = ":::notes:::"
NOTES_END = ":::endnotes:::"
IMAGE_START = ":::image:::"
IMAGE_END = ":::endimage:::"


def parse_markdown(md_text: str):
    """Parse Markdown into a list of slide dicts.

    Each slide: {title, layout_idx, body_lines, notes_lines, image_descs}
    """
    slides = []
    current = None
    in_notes = False
    in_image = False

    for raw in md_text.splitlines():
        line = raw.rstrip()

        m = SLIDE_HEADING_RE.match(line)
        if m and not in_notes and not in_image:
            if current is not None:
                slides.append(current)
            layout_idx = int(m.group(1)) if m.group(1) else 1
            title = m.group(2).strip()
            current = {
                "title": title,
                "layout_idx": layout_idx,
                "body_lines": [],
                "notes_lines": [],
                "image_descs": [],
            }
            continue

        if current is None:
            continue

        if line.strip() == NOTES_START:
            in_notes = True
            continue
        if line.strip() == NOTES_END:
            in_notes = False
            continue
        if line.strip() == IMAGE_START:
            in_image = True
            continue
        if line.strip() == IMAGE_END:
            in_image = False
            continue

        if in_notes:
            current["notes_lines"].append(line)
        elif in_image:
            current["image_descs"].append(line.strip())
        else:
            current["body_lines"].append(line)

    if current is not None:
        slides.append(current)

    return slides


def clean_body(lines):
    """Strip leading/trailing blank lines, normalize."""
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    return lines


def fill_text_frame(tf, body_lines):
    """Fill a text_frame with body lines. Bullets if line starts with `- `."""
    body_lines = clean_body(body_lines)
    if not body_lines:
        return

    tf.clear()
    first = True
    for line in body_lines:
        stripped = line.lstrip()
        indent_spaces = len(line) - len(stripped)
        is_bullet = stripped.startswith("- ")
        text = stripped[2:].strip() if is_bullet else stripped

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = text
        if is_bullet:
            level = min(indent_spaces // 2, 4)
            p.level = level


def build_slide(prs, slide_data, default_layout=1):
    """Add one slide to prs."""
    layout_idx = slide_data["layout_idx"]
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = default_layout
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    # Title
    title_set = False
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0 and shape.has_text_frame:
            shape.text_frame.text = slide_data["title"]
            title_set = True
            break

    # Body
    body_lines = clean_body(slide_data["body_lines"])
    if body_lines:
        body_filled = False
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                continue
            if shape.has_text_frame:
                fill_text_frame(shape.text_frame, body_lines)
                body_filled = True
                break

        # Fallback: no body placeholder found — insert a textbox
        if not body_filled:
            from pptx.util import Inches
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            height = prs.slide_height - Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            fill_text_frame(txBox.text_frame, body_lines)

    # Append image hints to the slide notes
    if slide_data["image_descs"]:
        slide_data["notes_lines"].append("")
        slide_data["notes_lines"].append("IMAGES STILL TO INSERT:")
        for d in slide_data["image_descs"]:
            slide_data["notes_lines"].append(f"  - {d}")

    # Speaker notes
    notes_lines = clean_body(slide_data["notes_lines"])
    if notes_lines:
        notes_text = "\n".join(notes_lines)
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def main():
    ap = argparse.ArgumentParser(description="Build PPTX from Markdown + Template")
    ap.add_argument("markdown", help="Markdown input file using the slide convention")
    ap.add_argument("--template", required=True, help="Template PPTX (layout master)")
    ap.add_argument("--output", required=True, help="Target PPTX")
    ap.add_argument("--default-layout", type=int, default=1, help="Default layout index")
    ap.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Keep the template's slides (default: the template is used only as a layout source and all of its existing slides are removed)",
    )
    args = ap.parse_args()

    md_path = Path(args.markdown)
    template_path = Path(args.template)
    output_path = Path(args.output)

    if not md_path.exists():
        sys.exit(f"Markdown not found: {md_path}")
    if not template_path.exists():
        sys.exit(f"Template not found: {template_path}")

    md_text = md_path.read_text(encoding="utf-8")
    slides = parse_markdown(md_text)
    if not slides:
        sys.exit("No slides found in the Markdown. Convention: '## Slide N: Title'")

    print(f"Markdown: {md_path.name}")
    print(f"Template: {template_path.name}")
    print(f"{len(slides)} slides found")

    prs = Presentation(str(template_path))

    if not args.keep_template_slides:
        # Remove all template slides; layouts/master stay
        slide_id_lst = prs.slides._sldIdLst
        old_count = len(slide_id_lst)
        for sld in list(slide_id_lst):
            slide_id_lst.remove(sld)
        print(f"Removed {old_count} template slides (layouts kept)")

    for i, slide_data in enumerate(slides, 1):
        build_slide(prs, slide_data, default_layout=args.default_layout)
        print(f"  Slide {i}: {slide_data['title'][:60]}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
